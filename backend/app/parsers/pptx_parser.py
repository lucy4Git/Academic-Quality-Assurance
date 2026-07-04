"""PPTX parser using python-pptx.

Extracts text from all shapes on all slides.  Title-type shapes (placeholder
type PP_PLACEHOLDER.TITLE / CENTER_TITLE) are collected separately as slide
titles in metadata, since they are the most content-dense signal for
classification.
"""

from __future__ import annotations

import io

from app.parsers.base import DocumentParser, ExtractionResult, ParserError


class PptxParser(DocumentParser):
    """Extract text from PPTX presentations."""

    @property
    def supported_mime_types(self) -> frozenset[str]:
        return frozenset(
            {
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            }
        )

    async def extract(self, content: bytes, filename: str) -> ExtractionResult:
        return await self._run_sync(self._extract_sync, content, filename)

    @staticmethod
    def _extract_sync(content: bytes, filename: str) -> ExtractionResult:
        try:
            from pptx import Presentation  # noqa: PLC0415
            from pptx.enum.shapes import PP_PLACEHOLDER  # noqa: PLC0415
        except ImportError:
            raise ParserError(
                "python-pptx is not installed. Run: pip install python-pptx",
                recoverable=False,
            )

        try:
            prs = Presentation(io.BytesIO(content))
        except Exception as exc:
            raise ParserError(f"Cannot open PPTX '{filename}': {exc}") from exc

        parts: list[str] = []
        slide_titles: list[str] = []

        title_placeholder_types = {
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
        }

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                shape_text = shape.text_frame.text.strip()
                if not shape_text:
                    continue
                parts.append(shape_text)

                # Collect title-type placeholders separately.
                if (
                    shape.is_placeholder
                    and hasattr(shape, "placeholder_format")
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.type in title_placeholder_types
                ):
                    slide_titles.append(shape_text)

        full_text = "\n".join(parts).strip()
        slide_count = len(prs.slides)

        return ExtractionResult(
            text=full_text,
            word_count=len(full_text.split()),
            parser_name="pptx",
            page_count=slide_count,
            metadata={
                "slide_count": slide_count,
                "slide_titles": slide_titles,
            },
        )
